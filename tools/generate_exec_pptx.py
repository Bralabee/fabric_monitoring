#!/usr/bin/env python3
"""Generate an executive-ready PowerPoint deck for USF Fabric Monitoring.

This script creates a 12-slide .pptx with a consistent, professional layout and
content grounded in the repository docs (README/WIKI/PROJECT_ANALYSIS).

Usage:
  python tools/generate_exec_pptx.py \
    --out docs/executive/USF_Fabric_Monitoring_Executive_Overview.pptx
"""

from __future__ import annotations

import argparse
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_NAME = "USF Fabric Monitoring & Governance"
DECK_SUBTITLE = "Executive overview (Microsoft Fabric)"


def _set_run_font(run, *, size: int | None = None, bold: bool | None = None, color: RGBColor | None = None):
    font = run.font
    if size is not None:
        font.size = Pt(size)
    if bold is not None:
        font.bold = bold
    if color is not None:
        font.color.rgb = color


def _set_shape_text(shape, title: str | None = None, bullets: list[str] | None = None, *, font_size: int = 20):
    tf = shape.text_frame
    tf.clear()

    if title:
        p = tf.paragraphs[0]
        p.text = title
        _set_run_font(p.runs[0], size=font_size, bold=True)

    if bullets:
        for i, bullet in enumerate(bullets):
            p = tf.add_paragraph() if (title or i > 0) else tf.paragraphs[0]
            p.text = bullet
            p.level = 0
            p.space_before = Pt(2)
            p.space_after = Pt(2)
            for r in p.runs:
                _set_run_font(r, size=16)


def _add_footer(slide, *, left_text: str, right_text: str):
    # Footer band
    left = Inches(0.5)
    top = Inches(7.0)
    width = Inches(12.33)
    height = Inches(0.35)
    band = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE = 1
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(245, 245, 245)
    band.line.color.rgb = RGBColor(230, 230, 230)

    # Left footer
    txl = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(8.0), Inches(0.25))
    p = txl.text_frame.paragraphs[0]
    p.text = left_text
    p.alignment = PP_ALIGN.LEFT
    _set_run_font(p.runs[0], size=10, color=RGBColor(80, 80, 80))

    # Right footer
    txr = slide.shapes.add_textbox(Inches(9.2), Inches(7.05), Inches(3.6), Inches(0.25))
    p = txr.text_frame.paragraphs[0]
    p.text = right_text
    p.alignment = PP_ALIGN.RIGHT
    _set_run_font(p.runs[0], size=10, color=RGBColor(80, 80, 80))


def _title_slide(prs: Presentation, *, subtitle: str, context: str):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = PROJECT_NAME

    st = slide.placeholders[1].text_frame
    st.clear()
    p1 = st.paragraphs[0]
    p1.text = subtitle
    _set_run_font(p1.runs[0], size=24)

    p2 = st.add_paragraph()
    p2.text = context
    _set_run_font(p2.runs[0], size=14, color=RGBColor(90, 90, 90))

    _add_footer(slide, left_text=DECK_SUBTITLE, right_text=str(date.today()))


def _bullets_slide(prs: Presentation, *, title: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        for r in p.runs:
            _set_run_font(r, size=18)

    _add_footer(slide, left_text=PROJECT_NAME, right_text=str(date.today()))


def _two_column_slide(prs: Presentation, *, title: str, left_title: str, left_bullets: list[str], right_title: str, right_bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    slide.shapes.title.text = title

    left_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.55), Inches(6.1), Inches(5.2))
    _set_shape_text(left_box, title=left_title, bullets=left_bullets, font_size=18)

    right_box = slide.shapes.add_textbox(Inches(6.9), Inches(1.55), Inches(6.1), Inches(5.2))
    _set_shape_text(right_box, title=right_title, bullets=right_bullets, font_size=18)

    _add_footer(slide, left_text=PROJECT_NAME, right_text=str(date.today()))


def _flow_slide(prs: Presentation, *, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    slide.shapes.title.text = title

    # Simple flow diagram using labeled boxes
    top = Inches(2.05)
    box_w = Inches(2.35)
    box_h = Inches(1.0)
    gap = Inches(0.35)
    left0 = Inches(0.7)

    labels = [
        "Fabric / Power BI APIs\n(Activity + Admin)",
        "Extract\n(per-day, MAX_HISTORICAL_DAYS)",
        "Enrich\n(job details ~8h cache)",
        "Reports\n(CSV + Parquet)",
        "Decisions\nOps + Governance",
    ]

    for idx, label in enumerate(labels):
        left = left0 + idx * (box_w + gap)
        box = slide.shapes.add_shape(1, left, top, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(250, 250, 250)
        box.line.color.rgb = RGBColor(200, 200, 200)
        tf = box.text_frame
        tf.text = label
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        for r in tf.paragraphs[0].runs:
            _set_run_font(r, size=14, bold=True, color=RGBColor(40, 40, 40))

        # Arrows (except last)
        if idx < len(labels) - 1:
            arrow_left = left + box_w
            arrow = slide.shapes.add_shape(13, arrow_left, top + Inches(0.35), gap, Inches(0.3))  # RIGHT_ARROW = 13
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(220, 220, 220)
            arrow.line.color.rgb = RGBColor(220, 220, 220)

    note = slide.shapes.add_textbox(Inches(0.7), Inches(3.45), Inches(12.6), Inches(2.7))
    _set_shape_text(
        note,
        title="Key design choices",
        bullets=[
            "Tenant-wide mode uses Power BI Admin Activity Events API (1 call/day) with automatic fallback to member-only scope when needed.",
            "Outputs are environment-agnostic: in Fabric they land under OneLake (`/lakehouse/default/Files/...`) via `resolve_path()`.",
            "Artifacts are designed for downstream consumption (Power BI dashboards, incident response, audit evidence).",
        ],
        font_size=16,
    )

    _add_footer(slide, left_text=PROJECT_NAME, right_text=str(date.today()))


def build_deck() -> Presentation:
    prs = Presentation()

    _title_slide(
        prs,
        subtitle=DECK_SUBTITLE,
        context="Purpose: provide operational visibility, governance controls, and evidence for Microsoft Fabric at tenant scale.",
    )

    _bullets_slide(
        prs,
        title="1. Executive Summary",
        bullets=[
            "Library-first monitoring and governance toolkit for Microsoft Fabric.",
            "Enables tenant-wide activity visibility (historical window capped by API limits) with consistent export artifacts.",
            "Adds governance guardrails: validates and enforces mandatory security group access across workspaces.",
            "Designed to run locally (Conda) and in Fabric (Environment + notebooks), producing repeatable evidence.",
        ],
    )

    _bullets_slide(
        prs,
        title="2. The Problem (What We’re Solving)",
        bullets=[
            "Limited centralized visibility into Fabric activity across all workspaces.",
            "Operational questions are hard to answer quickly: failures, long-running jobs, usage patterns, and drift.",
            "Access drift risk: required admin groups can be removed or downgraded over time.",
            "Need repeatable artifacts for governance reviews and executive reporting.",
        ],
    )

    _bullets_slide(
        prs,
        title="3. Solution Overview (Three Capabilities)",
        bullets=[
            "Monitor Hub: extract + analyze activity events and generate actionable reports.",
            "Access enforcement: assess/enforce mandatory security group roles with suppressions for labs/sandboxes.",
            "Lineage inventory: identify Mirrored Databases and capture source details for governance inventory.",
        ],
    )

    _flow_slide(prs, title="4. How It Works (High-Level Data Flow)")

    _two_column_slide(
        prs,
        title="5. Tenant-Wide Coverage Strategy",
        left_title="Tenant-wide (default)",
        left_bullets=[
            "Uses Power BI Admin Activity Events API to cover shared workspaces at scale.",
            "Operationally efficient: ~1 API call per day for activity events.",
            "Excludes personal workspaces to reduce noise and avoid known SP limitations.",
        ],
        right_title="Member-only (fallback / targeted)",
        right_bullets=[
            "Uses per-workspace enumeration for scenarios without admin permissions.",
            "Useful for development/testing or partial visibility scenarios.",
            "Selectable via `--member-only` (or `MEMBER_ONLY=1`).",
        ],
    )

    _bullets_slide(
        prs,
        title="6. Governance: Access Enforcement",
        bullets=[
            "Daily verification that mandatory Azure AD security groups retain expected roles in each workspace.",
            "Two-step operational guardrail: `assess` (dry run) → `enforce --confirm` (apply changes).",
            "Configuration-driven targets/suppressions via `config/workspace_access_targets.json` and `config/workspace_access_suppressions.json`.",
            "Produces timestamped JSON (and optional CSV) artifacts suitable for audit evidence and alerting.",
        ],
    )

    _bullets_slide(
        prs,
        title="7. Outputs & Executive KPIs (Examples)",
        bullets=[
            "Operational reporting: activity volume, success rate, top failing items/errors, long-running operations.",
            "Governance reporting: compliance rate by workspace, remediation actions required/applied, suppressed scope.",
            "Data products: CSV reports for broad consumption + Parquet “source of truth” for downstream analytics.",
            "All outputs default to `EXPORT_DIRECTORY=exports/monitor_hub_analysis` (OneLake-resolved in Fabric).",
        ],
    )

    _two_column_slide(
        prs,
        title="8. Deployment & Operations (How Teams Run It)",
        left_title="Local / CI runner",
        left_bullets=[
            "`make create` → `make install` → `make test-smoke`",
            "`make monitor-hub DAYS=7`",
            "`make enforce-access MODE=assess CSV_SUMMARY=1`",
        ],
        right_title="Fabric-native",
        right_bullets=[
            "Upload wheel to Fabric Environment and publish.",
            "Attach environment to notebooks under `notebooks/`.",
            "Artifacts land in OneLake under `/lakehouse/default/Files/exports/...`.",
        ],
    )

    _bullets_slide(
        prs,
        title="9. Roadmap (Pragmatic Next Steps)",
        bullets=[
            "Short-term: expand tests (especially enrichment/merge logic), tighten config validation, reduce notebook duplication.",
            "Medium-term: automated scheduling + alerting, dashboards/visualizations, richer lineage (pipelines/dataflows/models).",
            "Long-term: anomaly detection, enterprise audit trail for enforcement actions, broader observability integrations.",
        ],
    )

    _bullets_slide(
        prs,
        title="10. Risks / Dependencies",
        bullets=[
            "Permissions: tenant-wide monitoring relies on admin API permissions for the service principal.",
            "API limits: activity history window is capped by platform limits (configured via `MAX_HISTORICAL_DAYS`).",
            "Operational hygiene: define retention + ownership for artifacts under `exports/` and logs under `logs/`.",
            "Change management: enforcement requires explicit confirmation to avoid unintended access changes.",
        ],
    )

    _bullets_slide(
        prs,
        title="11. Executive Ask",
        bullets=[
            "Confirm governance goals: required groups/roles and the workspace suppression policy.",
            "Authorize tenant-wide permissions for the service principal (or approve member-only scope where necessary).",
            "Decide operating cadence (daily/weekly) and define owners for triage + remediation workflows.",
            "Greenlight a dashboard/reporting consumer (Power BI) for ongoing executive visibility.",
        ],
    )

    return prs


def main() -> int:
    parser = argparse.ArgumentParser(description="Generate an executive PowerPoint deck for USF Fabric Monitoring")
    parser.add_argument(
        "--out",
        type=str,
        default="docs/executive/USF_Fabric_Monitoring_Executive_Overview.pptx",
        help="Output .pptx path (relative to repo root by default)",
    )
    args = parser.parse_args()

    out_path = Path(args.out)
    if not out_path.is_absolute():
        out_path = Path.cwd() / out_path
    out_path.parent.mkdir(parents=True, exist_ok=True)

    prs = build_deck()
    prs.save(str(out_path))

    print(f"Wrote: {out_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
